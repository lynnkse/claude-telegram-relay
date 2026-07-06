#!/usr/bin/env python3
"""
Knowledge Base Ingestion Script

Extracts content from a PDF, slide deck, image, or link, chunks it, embeds via
Supabase embed function, stores in kb_items + kb_chunks.

Usage:
  python3 kb_ingest.py --pdf <path>   --title "..." [--author "..."] [--category "..."] [--source-url "..."]
  python3 kb_ingest.py --pptx <path>  --title "..." [--author "..."] [--category "..."] [--source-url "..."]
  python3 kb_ingest.py --image <path> --title "..." [--category "..."]
  python3 kb_ingest.py --link <url>   [--title "..."] [--category "..."]
"""

import os
import sys
import json
import time
import argparse
import subprocess
import urllib.request

SUPABASE_URL = os.environ.get("SUPABASE_URL", "")
SUPABASE_KEY = os.environ.get("SUPABASE_ANON_KEY", "")
CLAUDE_PATH = os.environ.get("CLAUDE_PATH", "claude")
CHUNK_CHARS = 1800
CHUNK_OVERLAP = 300


def _supabase_insert(table: str, row: dict) -> dict:
    data = json.dumps([row]).encode()
    req = urllib.request.Request(
        f"{SUPABASE_URL}/rest/v1/{table}",
        data=data,
        method="POST",
        headers={
            "apikey": SUPABASE_KEY,
            "Authorization": f"Bearer {SUPABASE_KEY}",
            "Content-Type": "application/json",
            "Prefer": "return=representation",
        },
    )
    with urllib.request.urlopen(req, timeout=30) as r:
        return json.loads(r.read().decode())[0]


def trigger_embed(row_id: str, content: str, table: str):
    data = json.dumps({"record": {"id": row_id, "content": content}, "table": table}).encode()
    req = urllib.request.Request(
        f"{SUPABASE_URL}/functions/v1/embed",
        data=data,
        headers={
            "Authorization": f"Bearer {SUPABASE_KEY}",
            "Content-Type": "application/json",
        },
    )
    with urllib.request.urlopen(req, timeout=30) as r:
        r.read()


def chunk_text(text: str) -> list[str]:
    chunks, start = [], 0
    text = text.strip()
    if not text:
        return []
    while start < len(text):
        end = min(start + CHUNK_CHARS, len(text))
        chunks.append(text[start:end])
        if end == len(text):
            break
        start += CHUNK_CHARS - CHUNK_OVERLAP
    return chunks


def create_item(title: str, author: str, category: str, source_url: str | None) -> str:
    row = _supabase_insert("kb_items", {
        "title": title,
        "author": author or None,
        "category": category or "general",
        "source_url": source_url,
    })
    return row["id"]


def ingest_chunks(item_id: str, text: str) -> tuple[int, int]:
    total, errors = 0, 0
    for idx, chunk in enumerate(chunk_text(text)):
        if len(chunk.strip()) < 30:
            continue
        try:
            row = _supabase_insert("kb_chunks", {
                "item_id": item_id,
                "content": chunk,
                "chunk_index": idx,
            })
            trigger_embed(row["id"], chunk, table="kb_chunks")
            total += 1
            time.sleep(0.1)
        except Exception as e:
            print(f"  ERROR chunk {idx}: {e}")
            errors += 1
    return total, errors


def extract_pdf_text(path: str) -> str:
    import pdfplumber
    parts = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            parts.append(page.extract_text() or "")
    return "\n".join(parts)


def extract_pptx_text(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        slide_text.append(line)
        if slide_text:
            parts.append(f"[Slide {i + 1}]\n" + "\n".join(slide_text))
    return "\n\n".join(parts)


def run_claude_oneshot(prompt: str, timeout: int = 120) -> str:
    result = subprocess.run(
        [CLAUDE_PATH, "--print", prompt],
        capture_output=True,
        text=True,
        timeout=timeout,
    )
    if result.returncode != 0:
        raise RuntimeError(f"claude --print failed: {result.stderr[-500:]}")
    return result.stdout.strip()


def extract_image_text(path: str) -> str:
    prompt = (
        f"Read the image at {path}. Transcribe all readable text verbatim "
        "(slide content, diagram labels, captions). If there is no text, "
        "describe what the image shows in 2-3 sentences. Output only the "
        "transcription/description, no commentary."
    )
    return run_claude_oneshot(prompt)


def extract_link_text(url: str) -> str:
    prompt = (
        f"Fetch {url} and extract the main article/page text content "
        "(skip navigation, ads, footers). Output only the extracted text, "
        "no commentary, no markdown formatting."
    )
    return run_claude_oneshot(prompt)


def main():
    parser = argparse.ArgumentParser()
    src = parser.add_mutually_exclusive_group(required=True)
    src.add_argument("--pdf")
    src.add_argument("--pptx")
    src.add_argument("--image")
    src.add_argument("--link")
    parser.add_argument("--title")
    parser.add_argument("--author", default="")
    parser.add_argument("--category", default="general")
    parser.add_argument("--source-url", default=None)
    args = parser.parse_args()

    if not SUPABASE_URL or not SUPABASE_KEY:
        print("Missing SUPABASE_URL or SUPABASE_ANON_KEY")
        sys.exit(1)

    if args.pdf:
        path = args.pdf
        title = args.title or os.path.splitext(os.path.basename(path))[0]
        print(f"Extracting PDF text: {path}")
        text = extract_pdf_text(path)
        source_url = args.source_url
    elif args.pptx:
        path = args.pptx
        title = args.title or os.path.splitext(os.path.basename(path))[0]
        print(f"Extracting PPTX text: {path}")
        text = extract_pptx_text(path)
        source_url = args.source_url
    elif args.image:
        path = args.image
        title = args.title or os.path.splitext(os.path.basename(path))[0]
        print(f"OCR/describing image via claude --print: {path}")
        text = extract_image_text(path)
        source_url = args.source_url
    else:
        url = args.link
        title = args.title or url
        print(f"Fetching link content via claude --print: {url}")
        text = extract_link_text(url)
        source_url = url

    if not text.strip():
        print("No content extracted — aborting, nothing written to knowledge base.")
        sys.exit(1)

    print(f"Extracted {len(text)} chars. Creating kb_items row...")
    item_id = create_item(title, args.author, args.category, source_url)
    print(f"kb_items id: {item_id}")

    total, errors = ingest_chunks(item_id, text)
    summary = f"Knowledge base ingestion complete!\nTitle: {title}\nCategory: {args.category}\nChunks: {total}\nErrors: {errors}"
    print(f"\n{summary}")


if __name__ == "__main__":
    main()
